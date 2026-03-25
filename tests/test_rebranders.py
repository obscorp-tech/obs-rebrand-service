"""Tests for DOCX and PPTX rebranding engines."""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

from rebrand_service.docx_rebrander import DocxRebrander
from rebrand_service.models import BrandConfig, ColorPalette, Typography
from rebrand_service.pptx_rebrander import PptxRebrander


@pytest.fixture
def brand() -> BrandConfig:
    return BrandConfig(
        client_name="Test Corp",
        client_slug="test-corp",
        colors=ColorPalette(
            primary="1B5E20",
            secondary="4CAF50",
            accent="FF9800",
            heading_text="1B5E20",
            body_text="212121",
        ),
        typography=Typography(
            heading_font="Georgia",
            body_font="Calibri",
            heading_size_pt=16,
            body_size_pt=11,
        ),
    )


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    """Create a minimal DOCX for testing."""
    doc = Document()
    doc.add_heading("Test Heading", level=1)
    doc.add_paragraph("Body text paragraph.")
    doc.add_heading("Sub Heading", level=2)
    doc.add_paragraph("Another paragraph with content.")
    path = tmp_path / "sample.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX for testing."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title and title.has_text_frame:
        title.text = "Test Title"
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body and body.has_text_frame:
        body.text = "Test body content."
    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


class TestDocxRebrander:
    def test_rebrand_produces_output(
        self, brand: BrandConfig, sample_docx: Path, tmp_path: Path
    ) -> None:
        output = tmp_path / "output" / "sample.docx"
        rebrander = DocxRebrander(brand, repo_root=tmp_path)
        result = rebrander.rebrand(sample_docx, output)

        assert result["status"] == "success"
        assert output.exists()
        assert result["input_sha256"] != result["output_sha256"]

    def test_typography_applied(
        self, brand: BrandConfig, sample_docx: Path, tmp_path: Path
    ) -> None:
        output = tmp_path / "output" / "sample.docx"
        rebrander = DocxRebrander(brand, repo_root=tmp_path)
        rebrander.rebrand(sample_docx, output)

        doc = Document(str(output))
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if paragraph.style.name.startswith("Heading"):
                    assert run.font.name == "Georgia"
                else:
                    assert run.font.name == "Calibri"

    def test_idempotent_double_run(
        self, brand: BrandConfig, sample_docx: Path, tmp_path: Path
    ) -> None:
        """Running rebrand twice should produce consistent results."""
        output1 = tmp_path / "out1" / "sample.docx"
        output2 = tmp_path / "out2" / "sample.docx"
        rebrander = DocxRebrander(brand, repo_root=tmp_path)

        rebrander.rebrand(sample_docx, output1)
        rebrander.rebrand(output1, output2)

        # Both outputs should have same fonts applied
        doc1 = Document(str(output1))
        doc2 = Document(str(output2))
        for p1, p2 in zip(doc1.paragraphs, doc2.paragraphs, strict=False):
            for r1, r2 in zip(p1.runs, p2.runs, strict=False):
                assert r1.font.name == r2.font.name

    def test_error_on_invalid_input(self, brand: BrandConfig, tmp_path: Path) -> None:
        fake_input = tmp_path / "nonexistent.docx"
        output = tmp_path / "output" / "nonexistent.docx"
        rebrander = DocxRebrander(brand, repo_root=tmp_path)
        result = rebrander.rebrand(fake_input, output)
        assert result["status"] == "error"

    def test_audit_hashes_populated(
        self, brand: BrandConfig, sample_docx: Path, tmp_path: Path
    ) -> None:
        output = tmp_path / "output" / "sample.docx"
        rebrander = DocxRebrander(brand, repo_root=tmp_path)
        result = rebrander.rebrand(sample_docx, output)

        assert len(result["input_sha256"]) == 64
        assert len(result["output_sha256"]) == 64


class TestPptxRebrander:
    def test_rebrand_produces_output(
        self, brand: BrandConfig, sample_pptx: Path, tmp_path: Path
    ) -> None:
        output = tmp_path / "output" / "sample.pptx"
        rebrander = PptxRebrander(brand, repo_root=tmp_path)
        result = rebrander.rebrand(sample_pptx, output)

        assert result["status"] == "success"
        assert output.exists()

    def test_error_on_invalid_input(self, brand: BrandConfig, tmp_path: Path) -> None:
        fake_input = tmp_path / "nonexistent.pptx"
        output = tmp_path / "output" / "nonexistent.pptx"
        rebrander = PptxRebrander(brand, repo_root=tmp_path)
        result = rebrander.rebrand(fake_input, output)
        assert result["status"] == "error"

    def test_compliance_footer_idempotent(
        self, brand: BrandConfig, sample_pptx: Path, tmp_path: Path
    ) -> None:
        """Compliance footer should not duplicate on re-run."""
        brand.compliance.confidentiality_label = "CONFIDENTIAL"
        brand.compliance.footer_text = "© Test Corp"

        output1 = tmp_path / "out1" / "sample.pptx"
        output2 = tmp_path / "out2" / "sample.pptx"

        rebrander = PptxRebrander(brand, repo_root=tmp_path)
        rebrander.rebrand(sample_pptx, output1)
        rebrander.rebrand(output1, output2)

        prs = Presentation(str(output2))
        footer_text = "CONFIDENTIAL | © Test Corp"
        for slide in prs.slides:
            count = sum(
                1
                for shape in slide.shapes
                if shape.has_text_frame
                and any(p.text.strip() == footer_text for p in shape.text_frame.paragraphs)
            )
            assert count <= 1, "Duplicate footer found on slide"
