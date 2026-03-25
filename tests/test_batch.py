"""Tests for batch processor."""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

from rebrand_service.batch import BatchProcessor, write_audit_log
from rebrand_service.models import BrandConfig, ColorPalette, Typography


@pytest.fixture
def brand() -> BrandConfig:
    return BrandConfig(
        client_name="Test Corp",
        client_slug="test-corp",
        colors=ColorPalette(primary="2E75B6", secondary="4A90D9", accent="F5A623"),
        typography=Typography(heading_font="Arial", body_font="Calibri"),
    )


@pytest.fixture
def input_dir_with_files(tmp_path: Path) -> Path:
    """Create directory with mixed test files."""
    input_dir = tmp_path / "input"
    input_dir.mkdir()

    # DOCX
    doc = Document()
    doc.add_paragraph("Test content")
    doc.save(str(input_dir / "report.docx"))

    # PPTX
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(input_dir / "deck.pptx"))

    # Unsupported file (should be skipped)
    (input_dir / "notes.txt").write_text("not a doc")

    # Temp file (should be skipped)
    doc2 = Document()
    doc2.save(str(input_dir / "~$temp.docx"))

    return input_dir


class TestBatchProcessor:
    def test_process_directory(
        self, brand: BrandConfig, input_dir_with_files: Path, tmp_path: Path
    ) -> None:
        output_dir = tmp_path / "output"
        processor = BatchProcessor(brand, repo_root=tmp_path)
        results = processor.process_directory(input_dir_with_files, output_dir)

        assert len(results) == 2  # docx + pptx (txt and ~$ skipped)
        assert all(r["status"] == "success" for r in results)
        assert (output_dir / "report.docx").exists()
        assert (output_dir / "deck.pptx").exists()

    def test_empty_directory(self, brand: BrandConfig, tmp_path: Path) -> None:
        empty_dir = tmp_path / "empty"
        empty_dir.mkdir()
        output_dir = tmp_path / "output"
        processor = BatchProcessor(brand, repo_root=tmp_path)
        results = processor.process_directory(empty_dir, output_dir)
        assert results == []


class TestAuditLog:
    def test_write_audit_log(self, tmp_path: Path) -> None:
        results = [
            {"input_file": "a.docx", "status": "success", "output_sha256": "abc123"},
            {"input_file": "b.pptx", "status": "error", "error": "bad file"},
        ]
        log_path = write_audit_log(results, tmp_path, "test-corp")

        assert log_path.exists()
        data = json.loads(log_path.read_text())
        assert data["total_files"] == 2
        assert data["success"] == 1
        assert data["errors"] == 1
        assert data["client"] == "test-corp"
