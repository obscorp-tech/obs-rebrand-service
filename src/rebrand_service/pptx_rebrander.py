"""PPTX rebranding engine — applies brand config to PowerPoint files."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from rebrand_service.models import BrandConfig, compute_file_hash

logger = logging.getLogger(__name__)


class PptxRebrander:
    """Applies brand configuration to PPTX files idempotently."""

    def __init__(self, brand: BrandConfig, repo_root: Path | None = None) -> None:
        self.brand = brand
        self.repo_root = repo_root or Path.cwd()

    def rebrand(self, input_path: Path, output_path: Path) -> dict[str, str]:
        """
        Rebrand a PPTX file and write to output_path.

        Returns audit metadata dict with input/output hashes.
        """
        try:
            input_hash = compute_file_hash(input_path)
            logger.info(
                "Rebranding PPTX: %s -> %s (client=%s, input_sha256=%s)",
                input_path.name,
                output_path.name,
                self.brand.client_slug,
                input_hash[:16],
            )
            prs = Presentation(str(input_path))
            self._apply_typography(prs)
            self._apply_colors(prs)
            self._apply_logo_to_slides(prs)
            self._apply_compliance_footer(prs)

            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))

            output_hash = compute_file_hash(output_path)
            logger.info(
                "Rebranded PPTX saved: %s (output_sha256=%s)",
                output_path.name,
                output_hash[:16],
            )

            return {
                "input_file": str(input_path),
                "output_file": str(output_path),
                "input_sha256": input_hash,
                "output_sha256": output_hash,
                "client": self.brand.client_slug,
                "status": "success",
            }

        except Exception as e:
            logger.exception("Failed to rebrand %s: %s", input_path.name, e)
            return {
                "input_file": str(input_path),
                "output_file": str(output_path),
                "client": self.brand.client_slug,
                "status": "error",
                "error": str(e),
            }

    def _apply_typography(self, prs: Presentation) -> None:
        """Apply font family and sizes to all text frames."""
        typo = self.brand.typography

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    is_title = self._is_title_shape(shape)
                    for run in paragraph.runs:
                        run.font.name = typo.heading_font if is_title else typo.body_font

    def _apply_colors(self, prs: Presentation) -> None:
        """Apply brand colors to all text."""
        colors = self.brand.colors
        RGBColor.from_string(colors.primary)
        body_rgb = RGBColor.from_string(colors.body_text)
        heading_rgb = RGBColor.from_string(colors.heading_text)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                is_title = self._is_title_shape(shape)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if is_title:
                            run.font.color.rgb = heading_rgb
                        else:
                            run.font.color.rgb = body_rgb

    def _apply_logo_to_slides(self, prs: Presentation) -> None:
        """Add logo to specified position on all slides."""
        if self.brand.logo is None:
            return

        logo_path = self.repo_root / self.brand.logo.path
        if not logo_path.exists():
            logger.warning("Logo file not found: %s", logo_path)
            return

        slide_width = prs.slide_width
        logo_width = Inches(self.brand.logo.width_inches)

        for slide in prs.slides:
            # Skip if logo already exists (idempotent)
            if self._slide_has_logo(slide):
                continue

            if self.brand.logo.position in ("header", "title-slide"):
                # Top-right corner
                left = slide_width - logo_width - Inches(0.5)
                top = Inches(0.3)
            else:
                # Footer position — bottom-right
                left = slide_width - logo_width - Inches(0.5)
                top = prs.slide_height - Inches(0.8)

            slide.shapes.add_picture(
                str(logo_path),
                left,
                top,
                width=logo_width,
            )

    def _apply_compliance_footer(self, prs: Presentation) -> None:
        """Add compliance text box to slide footers."""
        compliance = self.brand.compliance
        footer_parts: list[str] = []

        if compliance.confidentiality_label:
            footer_parts.append(compliance.confidentiality_label)
        if compliance.footer_text:
            footer_parts.append(compliance.footer_text)

        if not footer_parts:
            return

        footer_text = " | ".join(footer_parts)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide in prs.slides:
            # Check if compliance footer already exists (idempotent)
            if self._slide_has_compliance_footer(slide, footer_text):
                continue

            txbox = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=slide_height - Inches(0.5),
                width=slide_width - Inches(1.0),
                height=Inches(0.3),
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = footer_text
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor.from_string("999999")
            p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _is_title_shape(shape) -> bool:
        """Detect if a shape is a title placeholder."""
        if not shape.has_text_frame:
            return False
        try:
            ph = shape.placeholder_format
            return ph is not None and ph.idx in (0, 1)
        except ValueError:
            return False

    @staticmethod
    def _slide_has_logo(slide) -> bool:
        """Check if slide already has a picture (basic logo detection)."""
        from pptx.shapes.picture import Picture

        return any(isinstance(shape, Picture) for shape in slide.shapes)

    @staticmethod
    def _slide_has_compliance_footer(slide, footer_text: str) -> bool:
        """Check if compliance footer already exists on slide."""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip() == footer_text.strip():
                        return True
        return False
